#!/usr/bin/env python3
"""
Build a one-slide PowerPoint structure chart from a PSC tree (output of
walk_psc.py) in RockCap house style.

Three layout patterns:
  - chain:     simple vertical chain (default if topology is single-PSC at every level)
  - branching: two or more SPVs at top converging to a common UBO chain
  - jv:        non-100% splits, multiple ultimate owners side-by-side

Usage:
    python build_chart.py --input tree.json --output chart.pptx \\
        --title "Land off Broadway Hill, Horton — Borrower Structure Chart" \\
        --subtitle "Prepared by RockCap   |   11/05/2026" \\
        [--layout chain|branching|jv|auto]    (default auto)

The tree JSON shape (from walk_psc.py):
    {
        "company_number": "13614094",
        "name": "GALION (HORTON) LIMITED",
        ...
        "pscs": [
            {
                "kind": "corporate-entity-...",
                "name": "Galion Land Limited",
                "share_band": "75-100",
                "parent": { ... recursive same shape ... }
            },
            { "kind": "individual-...", "name": "Mr Macdonald", ... }
        ]
    }
"""
from __future__ import annotations
import argparse, json, os, sys
from dataclasses import dataclass
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# House-style colours (from references/style_guide.md)
NAVY  = RGBColor(0x0A, 0x2A, 0x4F)
LIGHT = RGBColor(0xD9, 0xE2, 0xEC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY  = RGBColor(0x4A, 0x4A, 0x4A)


# ---------------------------------------------------------------------------
# Tree → linear chain conversion
# ---------------------------------------------------------------------------
@dataclass
class ChainNode:
    title: str            # main label (uppercase)
    subtitle: str         # secondary label (e.g. "Co. 13614094" or "Ultimate beneficial owner")
    edge_label: str = ""  # label on the edge ABOVE this node (e.g. "100%")


def linearise(tree: dict) -> List[ChainNode]:
    """Walk a single-PSC chain into an ordered list of ChainNodes (top → bottom)."""
    nodes = [ChainNode(title=tree['name'].upper(),
                       subtitle=f"Co. {tree['company_number']}")]
    cursor = tree
    while cursor.get('pscs'):
        # Pick the corporate parent if there is one (preferred for chain layout);
        # otherwise the individual UBO terminates the chain.
        corp = next((p for p in cursor['pscs'] if p.get('is_corporate') and p.get('parent')), None)
        indiv = next((p for p in cursor['pscs'] if p.get('is_individual')), None)
        if corp:
            edge = _format_band(corp.get('share_band'))
            parent = corp['parent']
            nodes.append(ChainNode(
                title=parent['name'].upper(),
                subtitle=f"Co. {parent['company_number']}",
                edge_label=edge,
            ))
            cursor = parent
        elif indiv:
            edge = _format_band(indiv.get('share_band'))
            nodes.append(ChainNode(
                title=indiv['name'].upper(),
                subtitle='Ultimate beneficial owner',
                edge_label=edge,
            ))
            break
        else:
            break
    return nodes


def _format_band(band: Optional[str]) -> str:
    if not band:
        return ''
    # CH bands like '75-100' → render as '75-100%'; '100' (rare) → '100%'.
    return f'{band}%'


def detect_topology(tree: dict) -> str:
    """Heuristic: pick chain / branching / jv based on tree shape."""
    # If any node has >1 corporate PSC, treat as JV
    def has_multi_corp(node):
        corps = [p for p in node.get('pscs', []) if p.get('is_corporate')]
        if len(corps) > 1:
            return True
        for c in corps:
            if c.get('parent') and has_multi_corp(c['parent']):
                return True
        return False
    if has_multi_corp(tree):
        return 'jv'
    return 'chain'


# ---------------------------------------------------------------------------
# PPTX rendering
# ---------------------------------------------------------------------------
def _add_box(slide, left, top, width, height, title, subtitle):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid(); sh.fill.fore_color.rgb = LIGHT
    sh.line.color.rgb = NAVY; sh.line.width = Pt(1)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.10); tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.06); tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.name = 'Calibri'; r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = NAVY
    if subtitle:
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = subtitle
        r.font.name = 'Calibri'; r.font.size = Pt(10); r.font.color.rgb = NAVY


def _add_label(slide, left, top, width, height, text):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tb.fill.solid(); tb.fill.fore_color.rgb = WHITE
    tb.line.fill.background()
    tf = tb.text_frame
    tf.margin_left = Inches(0.03); tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01); tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = 'Calibri'; r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = NAVY


def _add_line(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = NAVY; conn.line.width = Pt(1)


def _add_title_block(slide, title, subtitle):
    t = slide.shapes.add_textbox(Inches(0.5), Inches(0.20), Inches(12.33), Inches(0.40))
    p = t.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.name = 'Calibri'; r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = NAVY
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(12.33), Inches(0.25))
        p = s.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = subtitle
        r.font.name = 'Calibri'; r.font.size = Pt(11); r.font.color.rgb = GREY


def _add_footer(slide, text):
    f = slide.shapes.add_textbox(Inches(0.5), Inches(7.22), Inches(12.33), Inches(0.25))
    p = f.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = text
    r.font.name = 'Calibri'; r.font.size = Pt(9); r.font.color.rgb = GREY


def render_chain(slide, nodes: List[ChainNode]):
    """Vertical chain layout, boxes centred horizontally."""
    box_w = 5.0
    left = (13.333 - box_w) / 2
    centre_x = 13.333 / 2
    n = len(nodes)
    if n == 0:
        return
    # Distribute vertically between y=1.10 and y=7.05
    top_band, bottom_band = 1.10, 7.05
    available = bottom_band - top_band
    # Each box ~0.95" tall (last box 0.75"), gap = remainder split evenly
    box_heights = [0.95] * (n - 1) + [0.75] if n >= 1 else []
    total_box_h = sum(box_heights)
    gap = (available - total_box_h) / max(n - 1, 1)
    if gap < 0.30:
        # Boxes don't fit; squeeze them down
        gap = 0.30
        scale = (available - gap * (n - 1)) / total_box_h
        box_heights = [h * scale for h in box_heights]

    y = top_band
    box_tops_bottoms = []
    for i, node in enumerate(nodes):
        h = box_heights[i]
        _add_box(slide, left, y, box_w, h, node.title, node.subtitle)
        box_tops_bottoms.append((y, y + h))
        y += h + gap

    # Connectors with labels
    for i in range(1, n):
        _, y_bot_prev = box_tops_bottoms[i - 1]
        y_top_curr, _ = box_tops_bottoms[i]
        _add_line(slide, centre_x, y_bot_prev, centre_x, y_top_curr)
        if nodes[i].edge_label:
            mid_y = (y_bot_prev + y_top_curr) / 2 - 0.11
            _add_label(slide, centre_x - 0.40, mid_y, 0.80, 0.22, nodes[i].edge_label)


def render_jv(slide, tree: dict):
    """Multi-PSC layout: render the SPV at top, then each PSC as a box on the next
    row, each pointing up to the SPV with its share-band label. Walks each PSC's
    parent chain as a sub-column.

    Kept deliberately minimal — escalate to a manual edit for genuinely complex
    JV structures rather than overengineering automated layout.
    """
    # SPV at top
    spv_w = 5.5
    spv_left = (13.333 - spv_w) / 2
    _add_box(slide, spv_left, 1.10, spv_w, 0.95,
             tree['name'].upper(), f"Co. {tree['company_number']}")
    spv_bottom = 2.05

    pscs = tree.get('pscs', [])
    n = len(pscs)
    if n == 0:
        return
    col_w = 13.333 / n
    box_w = min(col_w - 0.4, 4.5)

    for i, psc in enumerate(pscs):
        col_centre = (i + 0.5) * col_w
        col_left = col_centre - box_w / 2
        if psc.get('is_corporate'):
            _add_box(slide, col_left, 3.30, box_w, 0.95,
                     psc['name'].upper(),
                     f"Co. {psc.get('parent_company_number', '')}")
            # Recurse one level on parent if it exists
            parent = psc.get('parent')
            if parent:
                _add_box(slide, col_left, 4.95, box_w, 0.95,
                         parent['name'].upper(),
                         f"Co. {parent['company_number']}")
                _add_line(slide, col_centre, 4.25, col_centre, 4.95)
        else:
            # Individual UBO at the bottom of the column
            _add_box(slide, col_left, 3.30, box_w, 0.75,
                     psc['name'].upper(),
                     'Ultimate beneficial owner')
        # Connector up to SPV
        _add_line(slide, col_centre, 3.30, col_centre, spv_bottom)
        edge = _format_band(psc.get('share_band'))
        if edge:
            _add_label(slide, col_centre - 0.40, (3.30 + spv_bottom) / 2 - 0.11,
                       0.80, 0.22, edge)


# ---------------------------------------------------------------------------
def main():
    ap = argparse.ArgumentParser(description='Build a RockCap structure chart .pptx from a PSC tree.')
    ap.add_argument('--input', required=True, help='PSC tree JSON (output of walk_psc.py)')
    ap.add_argument('--output', required=True, help='Path to write .pptx')
    ap.add_argument('--title', required=True, help='Chart title (e.g. "Land off Broadway Hill, Horton — Borrower Structure Chart")')
    ap.add_argument('--subtitle', default='', help='Subtitle (default: "Prepared by RockCap   |   <today>")')
    ap.add_argument('--layout', default='auto', choices=['auto', 'chain', 'jv'])
    ap.add_argument('--footer', default='Source: Companies House PSC register.')
    args = ap.parse_args()

    with open(args.input, encoding='utf-8') as f:
        tree = json.load(f)

    layout = args.layout
    if layout == 'auto':
        layout = detect_topology(tree)

    if not args.subtitle:
        from datetime import date
        args.subtitle = f'Prepared by RockCap   |   {date.today().strftime("%d/%m/%Y")}'

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_block(slide, args.title, args.subtitle)

    if layout == 'chain':
        render_chain(slide, linearise(tree))
    elif layout == 'jv':
        render_jv(slide, tree)
    else:
        raise SystemExit(f'Unknown layout: {layout}')

    _add_footer(slide, args.footer)

    os.makedirs(os.path.dirname(os.path.abspath(args.output)) or '.', exist_ok=True)
    prs.save(args.output)
    print(f'Wrote: {args.output} ({layout} layout)')


if __name__ == '__main__':
    main()
